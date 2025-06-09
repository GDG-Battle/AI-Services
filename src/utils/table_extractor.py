import os
import pdfplumber
from docx import Document
from pptx import Presentation

def extract_tables_from_pdf(pdf_path, output_folder):
    os.makedirs(output_folder, exist_ok=True)
    all_tables_text = []

    with pdfplumber.open(pdf_path) as pdf:
        for page_num, page in enumerate(pdf.pages, start=1):
            tables = page.extract_tables()
            for table_index, table in enumerate(tables):
                # Filter out None values and convert to strings
                clean_table = [[str(cell).strip() if cell else "" for cell in row] for row in table]
                if clean_table:  # Only add tables that have content
                    table_text = f"[PDF Page {page_num} - Table {table_index+1}]\n"
                    table_text += "\n".join([" | ".join(row) for row in clean_table])
                    all_tables_text.append(table_text)

    # Save extracted tables to a text file
    if all_tables_text:
        table_file_path = os.path.join(output_folder, "tables.txt")
        with open(table_file_path, 'w', encoding='utf-8') as f:
            f.write("\n\n".join(all_tables_text))

def extract_tables_from_docx(docx_path, output_folder):
    os.makedirs(output_folder, exist_ok=True)
    doc = Document(docx_path)
    tables_text = []

    for table_index, table in enumerate(doc.tables):
        # Convert table to rows of strings
        table_rows = [[cell.text.strip() for cell in row.cells] for row in table.rows]
        table_text = f"\n[DOCX Table {table_index+1}]\n"
        table_text += "\n".join([" | ".join(row) for row in table_rows])
        tables_text.append(table_text)

    # Save extracted tables to a text file
    if tables_text:
        table_file_path = os.path.join(output_folder, "tables.txt")
        with open(table_file_path, 'w', encoding='utf-8') as f:
            f.write("\n\n".join(tables_text))

def extract_tables_from_pptx(pptx_path, output_folder):
    os.makedirs(output_folder, exist_ok=True)
    prs = Presentation(pptx_path)
    tables_text = []

    for slide_index, slide in enumerate(prs.slides):
        for shape_index, shape in enumerate(slide.shapes):
            # Check if shape is actually a table
            if hasattr(shape, 'has_table') and shape.has_table:
                table = shape.table
                # Verify minimum table structure (at least 2 rows or 2 columns)
                if len(table.rows) < 2 and len(table.columns) < 2:
                    continue
                    
                # Convert table to rows of strings
                table_rows = []
                for row in table.rows:
                    row_data = [cell.text.strip() for cell in row.cells]
                    table_rows.append(row_data)
                
                table_text = f"[PPTX Slide {slide_index+1} - Table {shape_index+1}]\n"
                table_text += "\n".join([" | ".join(row) for row in table_rows])
                tables_text.append(table_text)

    # Save extracted tables to a text file
    if tables_text:
        table_file_path = os.path.join(output_folder, "tables.txt")
        with open(table_file_path, 'w', encoding='utf-8') as f:
            f.write("\n\n".join(tables_text))
        
def extract_tables_by_format(path, output_folder):
    """
    Extract tables based on file format
    """
    ext = os.path.splitext(path)[-1].lower()
    if ext == ".pdf":
        return extract_tables_from_pdf(path, output_folder)
    elif ext == ".docx":
        return extract_tables_from_docx(path, output_folder)
    elif ext == ".pptx":
        return extract_tables_from_pptx(path, output_folder)
    else:
        print(f"[Warning] Unsupported file type for table extraction: {ext}")
