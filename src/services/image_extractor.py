import os
from PIL import Image
import fitz  # PyMuPDF
from io import BytesIO
import hashlib
from .image_processor import process_image_and_save

def extract_images_from_pdf(pdf_path, output_folder):
    os.makedirs(output_folder, exist_ok=True)
    output_file = os.path.join(output_folder, "image_analysis.txt")
    doc = fitz.open(pdf_path)
    seen_hashes = set()
    count = 0

    for page_index in range(len(doc)):
        for img_index, img in enumerate(doc.get_page_images(page_index)):
            xref = img[0]
            base_image = doc.extract_image(xref)
            image_bytes = base_image["image"]

            # Check image dimensions using PIL
            img_stream = BytesIO(image_bytes)
            pil_img = Image.open(img_stream).convert("RGB")
            width, height = pil_img.size

            # Skip tiny images
            if width <= 2 or height <= 2:
                continue

            hash_val = hashlib.sha256(image_bytes).hexdigest()
            if hash_val in seen_hashes:
                continue
            seen_hashes.add(hash_val)

            # Process image
            source_info = f"[PDF Page {page_index+1}, Image {img_index+1}]"
            process_image_and_save(pil_img, output_file, source_info)
            count += 1

    print(f"[PDF] Processed {count} images from '{os.path.basename(pdf_path)}'")

def extract_images_from_pptx(pptx_path, output_folder):
    import os
    from pptx import Presentation
    from PIL import Image
    from io import BytesIO
    import hashlib
    from .image_processor import process_image_and_save

    os.makedirs(output_folder, exist_ok=True)
    output_file = os.path.join(output_folder, "image_analysis.txt")
    prs = Presentation(pptx_path)
    seen_hashes = set()
    count = 0

    for slide_idx, slide in enumerate(prs.slides):
        for shape_idx, shape in enumerate(slide.shapes):
            if shape.shape_type == 13:  # PICTURE
                image = shape.image
                image_bytes = image.blob

                img_stream = BytesIO(image_bytes)
                pil_img = Image.open(img_stream).convert("RGB")
                width, height = pil_img.size

                if width <= 2 or height <= 2:
                    continue

                hash_val = hashlib.sha256(image_bytes).hexdigest()
                if hash_val in seen_hashes:
                    continue
                seen_hashes.add(hash_val)

                # Process image
                source_info = f"[PPTX Slide {slide_idx+1}, Image {shape_idx+1}]"
                process_image_and_save(pil_img, output_file, source_info)
                count += 1

    print(f"[PPTX] Processed {count} images from '{os.path.basename(pptx_path)}'")

def extract_images_from_docx(docx_path, output_folder):
    import os
    import docx2txt
    from PIL import Image
    from .image_processor import process_image_and_save
    
    os.makedirs(output_folder, exist_ok=True)
    temp_img_dir = os.path.join(output_folder, "temp_images")
    os.makedirs(temp_img_dir, exist_ok=True)
    
    # Extract images to temporary directory
    docx2txt.process(docx_path, temp_img_dir)
    
    output_file = os.path.join(output_folder, "image_analysis.txt")
    count = 0
    
    # Process each extracted image
    for img_idx, filename in enumerate(os.listdir(temp_img_dir)):
        if filename.lower().endswith(('.png', '.jpg', '.jpeg', '.gif')):
            img_path = os.path.join(temp_img_dir, filename)
            try:
                with Image.open(img_path).convert("RGB") as img:
                    source_info = f"[DOCX Image {img_idx+1}]"
                    process_image_and_save(img, output_file, source_info)
                count += 1
            except Exception as e:
                print(f"Error processing {filename}: {e}")
            finally:
                os.remove(img_path)  # Clean up temporary image
                
    os.rmdir(temp_img_dir)  # Remove temporary directory
    print(f"[DOCX] Processed {count} images from '{os.path.basename(docx_path)}'")