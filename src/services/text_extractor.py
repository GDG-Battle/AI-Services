import os

def extract_text_from_pdf(pdf_path, output_folder):
    import pdfplumber
    text_blocks = []
    with pdfplumber.open(pdf_path) as pdf:
        for page_num, page in enumerate(pdf.pages, start=1):
            raw_text = page.extract_text()
            if raw_text:
                page_text = f"\n[PDF Page {page_num}]\n{raw_text.strip()}"
                text_blocks.append(page_text)
    
    output_text = "\n".join(text_blocks)
    os.makedirs(output_folder, exist_ok=True)
    output_path = os.path.join(output_folder, "extracted_text.txt")
    with open(output_path, 'w', encoding='utf-8') as f:
        f.write(output_text)

def extract_text_from_docx(docx_path, output_folder):
    from docx import Document
    doc = Document(docx_path)
    text_blocks = []

    for para in doc.paragraphs:
        if para.text.strip():  # ignore empty lines
            text_blocks.append(para.text.strip())

    output_text = "\n".join(text_blocks)
    os.makedirs(output_folder, exist_ok=True)
    output_path = os.path.join(output_folder, "extracted_text.txt")
    with open(output_path, 'w', encoding='utf-8') as f:
        f.write(output_text)

def extract_text_from_pptx(pptx_path, output_folder):
    from pptx import Presentation
    prs = Presentation(pptx_path)
    text_blocks = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        slide_text = f"\n[PPTX Slide {slide_num}]\n"
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text.strip()
                if text:
                    slide_text += text + "\n"
        text_blocks.append(slide_text.strip())

    output_text = "\n".join(text_blocks)
    os.makedirs(output_folder, exist_ok=True)
    output_path = os.path.join(output_folder, "extracted_text.txt")
    with open(output_path, 'w', encoding='utf-8') as f:
        f.write(output_text)

def extract_text_by_format(path, output_folder):
    ext = os.path.splitext(path)[-1].lower()
    if ext == ".pdf":
        return extract_text_from_pdf(path, output_folder)
    elif ext == ".docx":
        return extract_text_from_docx(path, output_folder)
    elif ext == ".pptx":
        return extract_text_from_pptx(path, output_folder)
    else:
        return "[Unsupported file type]"